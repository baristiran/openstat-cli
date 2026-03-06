"""Export commands: export docx, export pptx."""

from __future__ import annotations

import os
from datetime import date
from pathlib import Path

from openstat.commands.base import command, CommandArgs
from openstat.session import Session


def _ensure_dir(path: str) -> None:
    os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)


# ── Word (.docx) ───────────────────────────────────────────────────────────

def _export_docx(session: Session, path: str) -> str:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        return (
            "python-docx is required for Word export.\n"
            "Install: pip install python-docx"
        )

    doc = Document()
    doc.add_heading("OpenStat Results", 0)
    doc.add_paragraph(
        f"Dataset: {session.dataset_name or 'Unknown'}  |  "
        f"Date: {date.today().isoformat()}  |  "
        f"Shape: {session.shape_str}"
    )
    doc.add_paragraph()

    # Dataset overview table
    doc.add_heading("Dataset Overview", level=1)
    if session.df is not None:
        df = session.df
        tbl = doc.add_table(rows=1, cols=2)
        tbl.style = "Table Grid"
        hdr = tbl.rows[0].cells
        hdr[0].text = "Property"
        hdr[1].text = "Value"
        for label, val in [
            ("Rows", str(df.height)),
            ("Columns", str(df.width)),
            ("Missing cells", str(sum(df[c].null_count() for c in df.columns))),
            ("Numeric columns", str(sum(1 for c in df.columns if df[c].dtype in (
                __import__("polars").Float32, __import__("polars").Float64,
                __import__("polars").Int32, __import__("polars").Int64,
            )))),
        ]:
            row = tbl.add_row().cells
            row[0].text = label
            row[1].text = val
        doc.add_paragraph()

        # Summary statistics table
        import polars as pl
        NUMERIC = (pl.Float32, pl.Float64, pl.Int8, pl.Int16, pl.Int32, pl.Int64,
                   pl.UInt8, pl.UInt16, pl.UInt32, pl.UInt64)
        num_cols = [c for c in df.columns if df[c].dtype in NUMERIC]
        if num_cols:
            doc.add_heading("Summary Statistics", level=2)
            stats_tbl = doc.add_table(rows=1, cols=5)
            stats_tbl.style = "Table Grid"
            for i, hdr_text in enumerate(["Variable", "N", "Mean", "SD", "Min–Max"]):
                stats_tbl.rows[0].cells[i].text = hdr_text
            for c in num_cols[:20]:  # cap at 20 rows
                col = df[c].drop_nulls()
                if col.len() == 0:
                    continue
                cells = stats_tbl.add_row().cells
                cells[0].text = c
                cells[1].text = str(col.len())
                cells[2].text = f"{col.mean():.4f}"
                cells[3].text = f"{col.std():.4f}" if col.len() > 1 else "—"
                cells[4].text = f"{col.min():.2f} – {col.max():.2f}"

    # Model results
    for mr in session.results:
        doc.add_heading(f"{mr.name}: {mr.formula}", level=1)
        doc.add_paragraph(mr.table, style="No Spacing")
        doc.add_paragraph()

    # Plots
    for plot_path in session.plot_paths:
        if os.path.exists(plot_path):
            doc.add_heading("Figure", level=2)
            try:
                doc.add_picture(plot_path)
            except Exception:
                doc.add_paragraph(f"[Plot: {plot_path}]")

    _ensure_dir(path)
    doc.save(path)
    return os.path.abspath(path)


# ── PowerPoint (.pptx) ────────────────────────────────────────────────────

def _export_pptx(session: Session, path: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return (
            "python-pptx is required for PowerPoint export.\n"
            "Install: pip install python-pptx"
        )

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank
    title_layout = prs.slide_layouts[0]  # title slide

    # Slide 1: Title
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "OpenStat Results"
    slide.placeholders[1].text = (
        f"{session.dataset_name or 'Dataset'}  |  {date.today().isoformat()}"
    )

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Overview"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = session.shape_str
    if session.df is not None:
        df = session.df
        tf.add_paragraph().text = f"Missing cells: {sum(df[c].null_count() for c in df.columns)}"

    # One slide per model
    for mr in session.results:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{mr.name}: {mr.formula}"
        tf = slide.placeholders[1].text_frame
        tf.text = mr.table[:1000]  # truncate if huge

    # One slide per plot
    for plot_path in session.plot_paths:
        if os.path.exists(plot_path):
            slide = prs.slides.add_slide(blank_layout)
            try:
                slide.shapes.add_picture(
                    plot_path,
                    Inches(0.5), Inches(0.5),
                    width=Inches(8), height=Inches(5.5),
                )
            except Exception:
                pass

    _ensure_dir(path)
    prs.save(path)
    return os.path.abspath(path)


# ── Command ────────────────────────────────────────────────────────────────

@command("export", usage="export docx|pptx|pdf|md [path]")
def cmd_export(session: Session, args: str) -> str:
    """Export results to Word (.docx), PowerPoint (.pptx), PDF, or Markdown."""
    ca = CommandArgs(args)
    if not ca.positional:
        return "Usage: export docx|pptx|pdf|md [path]"

    fmt = ca.positional[0].lower()

    if fmt == "docx":
        path = ca.positional[1] if len(ca.positional) > 1 else "outputs/results.docx"
        out = _export_docx(session, path)
        if out.endswith(".docx"):
            return f"Word document saved: {out}"
        return out

    elif fmt == "pptx":
        path = ca.positional[1] if len(ca.positional) > 1 else "outputs/results.pptx"
        out = _export_pptx(session, path)
        if out.endswith(".pptx"):
            return f"PowerPoint saved: {out}"
        return out

    elif fmt == "pdf":
        from openstat.commands.pdf_cmds import _export_pdf
        path = ca.positional[1] if len(ca.positional) > 1 else "outputs/results.pdf"
        out = _export_pdf(session, path)
        if out.endswith(".pdf"):
            return f"PDF saved: {out}"
        return out

    elif fmt == "md":
        from openstat.commands.pdf_cmds import _export_md
        path = ca.positional[1] if len(ca.positional) > 1 else "outputs/results.md"
        out = _export_md(session, path)
        return f"Markdown saved: {out}"

    else:
        return f"Unknown export format: {fmt}. Use 'docx', 'pptx', 'pdf', or 'md'."
